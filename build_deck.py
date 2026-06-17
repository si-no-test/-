#!/usr/bin/env python3
# -*- coding: utf-8 -*- 
"""
build_deck.py — 把整理好的新聞 (news_data.json) 產生成簡報。
v2 版型：無封面；每則一頁；分類標籤在右上角；標題與內文加大；日期顏色加深；
內文雙欄填滿版面（無資料表時左右兩欄；有資料表時左欄內文、右欄放表）。

用法:
    python build_deck.py news_data.json output.pptx
"""
import json
import sys
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SLIDE_W = 11949113
SLIDE_H = 6721475
EMU = 914400

FONT_ZH = "Microsoft JhengHei"
FONT_EN = "Calibri"

INK = RGBColor(0x26, 0x26, 0x26)
GREY = RGBColor(0x59, 0x59, 0x59)
BLUE = RGBColor(0x00, 0x70, 0xC0)
RED = RGBColor(0xC0, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROWALT = RGBColor(0xEE, 0xF3, 0xFA)

# 分類專屬色：邊框 / 盒底 / 標題文字
CAT_STYLE = {
    "金融政策與營運": (RGBColor(0xE0, 0x7A, 0x5F), RGBColor(0xFB, 0xE5, 0xD6), RGBColor(0xC0, 0x50, 0x4D)),
    "零售業務":       (RGBColor(0x70, 0xAD, 0x47), RGBColor(0xE2, 0xF0, 0xD9), RGBColor(0x54, 0x7C, 0x3A)),
    "法金/海外業務":   (RGBColor(0x80, 0x80, 0x80), RGBColor(0xF2, 0xF2, 0xF2), RGBColor(0x40, 0x40, 0x40)),
    "支付/金融科技":   (RGBColor(0x00, 0x70, 0xC0), RGBColor(0xD7, 0xDE, 0xED), RGBColor(0x1F, 0x49, 0x7D)),
}

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "assets", "sinopac_logo.png")


def set_run(run, text, size, color=INK, bold=False):
    run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT_EN
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", FONT_ZH)


def add_text(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Emu(int(l * EMU)), Emu(int(t * EMU)), Emu(int(w * EMU)), Emu(int(h * EMU)))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{side}", Emu(int(0.04 * EMU)))
    return tb, tf


def rounded(slide, l, t, w, h, fill, line):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(l * EMU)), Emu(int(t * EMU)), Emu(int(w * EMU)), Emu(int(h * EMU)))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    return sp


def footer(slide, page_no):
    tb, tf = add_text(slide, 0.4, 6.86, 8.0, 0.4)
    set_run(tf.paragraphs[0].add_run(), f"{page_no}  翻轉金融 共創美好生活 Together, a better life.", 11, GREY)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Emu(int(10.0 * EMU)), Emu(int(6.82 * EMU)), height=Emu(int(0.38 * EMU)))


def fill_paras(tf, paras, size, gap):
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = gap
        set_run(p.add_run(), para, size, INK)


def title_size(t):
    n = len(t)
    return 30 if n <= 20 else (26 if n <= 30 else 24)


def body_size(total, has_table):
    if has_table:
        ladder = [(1500, 12), (1150, 13), (850, 14), (600, 15), (0, 16)]
    else:
        ladder = [(2200, 12), (1700, 13), (1250, 14), (850, 15), (0, 16)]
    return next(sz for thr, sz in ladder if total >= thr)


def split_two(paras):
    total = sum(len(p) for p in paras)
    half = total * 0.55
    acc = 0; left = []; right = []
    for p in paras:
        if acc < half or not left:
            left.append(p); acc += len(p)
        else:
            right.append(p)
    return left, right


def build_detail(prs, art, page_no, date_range):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # 導讀小標
    tb, tf = add_text(s, 0.35, 0.12, 8.0, 0.32)
    set_run(tf.paragraphs[0].add_run(), f"【新聞導讀 {date_range}】", 11, GREY, bold=True)

    # 標題（加大）
    tb, tf = add_text(s, 0.4, 0.46, 8.0, 1.05)
    set_run(tf.paragraphs[0].add_run(), art["title"], title_size(art["title"]), INK, bold=True)

    # 右上角分類標籤
    cat = art.get("category", "")
    border, fillc, titlec = CAT_STYLE.get(cat, (GREY, RGBColor(0xF2, 0xF2, 0xF2), GREY))
    if cat:
        box = rounded(s, 10.2, 0.42, 2.5, 0.52, fillc, border)
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        set_run(pp.add_run(), cat, 14, titlec, bold=True)

    # 日期＋來源（加深、放大）
    tb, tf = add_text(s, 8.6, 1.12, 4.1, 0.4)
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    set_run(pp.add_run(), f"{art.get('date','')}　{art.get('source','')}", 14, INK, bold=True)

    summary = art.get("summary", [])
    total = sum(len(x) for x in summary)
    has_table = bool(art.get("table"))
    size = body_size(total, has_table)
    gap = Pt(6) if size <= 13 else Pt(8)

    body_top = 1.7
    body_h = 4.6

    if has_table:
        # 內文左欄（滿高）
        tb, tf = add_text(s, 0.4, body_top, 6.35, body_h)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        fill_paras(tf, summary, size, gap)
        # 右欄放資料表
        t = art["table"]
        tx, tw = 6.85, 5.85
        tb, tf = add_text(s, tx, 1.55, tw, 0.36)
        set_run(tf.paragraphs[0].add_run(), t.get("title", ""), 14, INK, bold=True)
        headers = t.get("headers", []); rows = t.get("rows", [])
        if headers and rows:
            nrow = len(rows) + 1; ncol = len(headers)
            th = min(3.9, 0.34 * nrow + 0.1)
            gf = s.shapes.add_table(nrow, ncol, Emu(int(tx * EMU)), Emu(int(2.0 * EMU)), Emu(int(tw * EMU)), Emu(int(th * EMU)))
            tbl = gf.table
            for j, htext in enumerate(headers):
                c = tbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = BLUE
                c.margin_top = Emu(int(0.02 * EMU)); c.margin_bottom = Emu(int(0.02 * EMU))
                pp = c.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
                set_run(pp.add_run(), str(htext), 11, WHITE, bold=True)
            for i, row in enumerate(rows, start=1):
                for j in range(ncol):
                    c = tbl.cell(i, j); c.fill.solid()
                    c.fill.fore_color.rgb = WHITE if i % 2 else ROWALT
                    c.margin_top = Emu(int(0.02 * EMU)); c.margin_bottom = Emu(int(0.02 * EMU))
                    set_run(c.text_frame.paragraphs[0].add_run(), str(row[j] if j < len(row) else ""), 10.5, INK)
            tb, tf = add_text(s, tx, 2.0 + th + 0.05, tw, 0.3)
            pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
            set_run(pp.add_run(), f"資料來源：{t.get('source','新聞整理')}", 9, GREY)
    else:
        # 無資料表：左右兩欄填滿
        left, right = split_two(summary)
        tb, tf = add_text(s, 0.4, body_top, 6.15, body_h)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        fill_paras(tf, left, size, gap)
        if right:
            tb, tf = add_text(s, 6.85, body_top, 5.85, body_h)
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            fill_paras(tf, right, size, gap)

    # 原文連結
    if art.get("url"):
        tb, tf = add_text(s, 0.4, 6.42, 8.5, 0.32)
        p = tf.paragraphs[0]
        set_run(p.add_run(), "原文連結：", 11, GREY)
        r = p.add_run(); set_run(r, art["url"], 11, BLUE); r.hyperlink.address = art["url"]
    footer(s, page_no)


def main():
    if len(sys.argv) < 3:
        print("用法: python build_deck.py news_data.json output.pptx"); sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        data = json.load(f)
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W); prs.slide_height = Emu(SLIDE_H)
    drange = data.get("date_range", "")
    arts = data.get("articles", [])
    for i, art in enumerate(arts, start=1):
        build_detail(prs, art, i, drange)
    prs.save(sys.argv[2])
    print(f"已產生 {sys.argv[2]}，共 {len(arts)} 頁（無封面）")


if __name__ == "__main__":
    main()
